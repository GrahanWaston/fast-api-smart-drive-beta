# routers/documents.py

from io import BytesIO
import json
import os
import secrets
import shutil
import tempfile
from fastapi import APIRouter, Depends, Query, Response, UploadFile, File, Form, HTTPException, BackgroundTasks
import pandas as pd
from pptx import Presentation
from sqlalchemy import or_, and_
from sqlalchemy.orm import Session, joinedload, selectinload, defer
from typing import List, Optional
from pdfminer.high_level import extract_text
from docx import Document as DocxDoc
from PIL import Image
import pytesseract
import openpyxl
from datetime import datetime, timedelta, timezone
from utils.authorization import *

from connection.database import SessionLocal
from models.models import Document, DocumentCategory, DocumentContent, StatusEnum, Directory, DocumentShare
from connection.schemas import (
    DocumentCategoryOut, DocumentCreate, DocumentOut, ContentOut, BulkActionRequest, DocumentShareCreate,
    DocumentShareOut, StatusUpdateResponse
)

router = APIRouter(prefix="/docs", tags=["documents"])


def get_db():
    db = SessionLocal()
    try:
        yield db
    finally:
        db.close()


def detect_file_category(mimetype: str) -> str:
    """Auto-detect file category from mimetype"""
    if not mimetype:
        return "Other"

    mimetype = mimetype.lower()

    if mimetype.startswith("image/"):
        return "Photo"
    elif "word" in mimetype or "msword" in mimetype:
        return "Word"
    elif "excel" in mimetype or "spreadsheet" in mimetype:
        return "Excel"
    elif "powerpoint" in mimetype or "presentation" in mimetype:
        return "PowerPoint"
    elif mimetype == "application/pdf":
        return "PDF"
    elif mimetype.startswith("video/"):
        return "Video"
    elif mimetype.startswith("audio/"):
        return "Audio"
    elif "zip" in mimetype or "compressed" in mimetype or "rar" in mimetype:
        return "Archive"
    else:
        return "Other"


@router.post("/", response_model=DocumentOut)
async def upload_document(
    background_tasks: BackgroundTasks,
    file: UploadFile = File(...),
    title_document: str = Form(None),
    directory_id: int = Form(None),

    # ✅ NEW PARAMETERS
    file_type: str = Form("Document"),  # Document or File
    document_category_id: int = Form(None),
    expire_date: str = Form(None),  # Format: YYYY-MM-DD

    # Existing parameters
    tags: str = Form(None),
    description: str = Form(None),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """
    Upload document with enhanced attributes

    New Fields:
    - file_type: Document or File
    - document_category_id: Category for this document (optional)
    - expire_date: Expiration date (YYYY-MM-DD format)

    Auto-filled:
    - file_category: Auto-detected from file type
    - file_owner: Current user's name
    - department_id: Current user's department
    - organization_id: Current user's organization
    """

    print(f"📤 UPLOAD DOCUMENT - User: {current_user.name}", flush=True)
    print(f"📤 File: {file.filename}, Type: {file_type}", flush=True)

    # Validate user has department and organization
    if not current_user.department_id or not current_user.organization_id:
        raise HTTPException(
            400, "User must be assigned to a department and organization")

    # Check directory access
    if directory_id:
        directory = db.query(Directory).filter(
            Directory.id == directory_id).first()
        if not directory:
            raise HTTPException(404, "Directory not found")
        if not can_access_directory(current_user, directory):
            raise HTTPException(403, "Access denied to this directory")

    # ✅ Validate document category if provided
    if document_category_id:
        doc_category = db.query(DocumentCategory).filter(
            DocumentCategory.id == document_category_id
        ).first()
        if not doc_category:
            raise HTTPException(404, "Document category not found")
        if doc_category.organization_id != current_user.organization_id:
            raise HTTPException(
                403, "Invalid document category for your organization")
        print(f"📁 Using category: {doc_category.name}", flush=True)

    # Read file data
    content_bytes = await file.read()
    size = len(content_bytes)

    # ✅ Auto-detect file category from mimetype
    file_category = detect_file_category(file.content_type)
    print(f"🔍 Detected file category: {file_category}", flush=True)

    # ✅ Parse expire_date
    parsed_expire_date = None
    if expire_date:
        try:
            parsed_expire_date = datetime.datetime.strptime(
                expire_date, "%Y-%m-%d")
            print(f"📅 Expire date set: {parsed_expire_date}", flush=True)
        except ValueError:
            raise HTTPException(
                400, "Invalid expire date format. Use YYYY-MM-DD")

    # ✅ Create document with new attributes
    doc = Document(
        name=file.filename,
        title_document=title_document or file.filename,

        # NEW ATTRIBUTES
        file_type=file_type,
        document_category_id=document_category_id,
        file_category=file_category,
        file_owner=current_user.name,  # Auto-fill
        expire_date=parsed_expire_date,

        # Existing attributes
        mimetype=file.content_type,
        size=size,
        data=content_bytes,
        directory_id=directory_id,
        status=StatusEnum.ACTIVE,
        organization_id=current_user.organization_id,  # Auto-fill
        department_id=current_user.department_id,  # Auto-fill
        created_by=current_user.id
    )

    db.add(doc)
    db.commit()
    db.refresh(doc)

    # Save metadata
    if tags or description:
        from models.models import DocumentMetadata
        metadata = DocumentMetadata(
            document_id=doc.id,
            tags=tags,
            description=description
        )
        db.add(metadata)
        db.commit()

    # Background text extraction
    # Di dalam upload_document function, ganti background task function


    def extract_and_store(doc_id: int, data: bytes, ext: str, mimetype: str):
        """Extract text from various file formats"""
        text = ""
        tmp_path = None

        try:
            # Create temporary file
            with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmpf:
                tmpf.write(data)
                tmpf.flush()
                tmp_path = tmpf.name

            # PDF extraction
            if ext.lower() == "pdf":
                try:
                    from pdfminer.high_level import extract_text
                    text = extract_text(tmp_path)
                    print(f"✅ Extracted {len(text)} chars from PDF", flush=True)
                except Exception as e:
                    print(f"❌ PDF extraction failed: {e}", flush=True)
                    text = ""

            # Word documents
            elif ext.lower() in ("docx", "doc"):
                try:
                    from docx import Document as DocxDoc
                    d = DocxDoc(tmp_path)
                    paragraphs = [p.text for p in d.paragraphs if p.text.strip()]
                    text = "\n".join(paragraphs)
                    print(f"✅ Extracted {len(text)} chars from Word", flush=True)
                except Exception as e:
                    print(f"❌ Word extraction failed: {e}", flush=True)
                    text = ""

            # Excel spreadsheets
            elif ext.lower() in ("xlsx", "xls"):
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(
                        tmp_path, read_only=True, data_only=True)
                    parts = []
                    for sheet in wb.worksheets:
                        for row in sheet.iter_rows(values_only=True):
                            row_text = " ".join(str(c)
                                                for c in row if c is not None)
                            if row_text.strip():
                                parts.append(row_text)
                    text = "\n".join(parts)
                    wb.close()
                    print(f"✅ Extracted {len(text)} chars from Excel", flush=True)
                except Exception as e:
                    print(f"❌ Excel extraction failed: {e}", flush=True)
                    text = ""

            # PowerPoint presentations
            elif ext.lower() in ("pptx", "ppt"):
                try:
                    from pptx import Presentation
                    prs = Presentation(tmp_path)
                    parts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                parts.append(shape.text)
                    text = "\n".join(parts)
                    print(
                        f"✅ Extracted {len(text)} chars from PowerPoint", flush=True)
                except Exception as e:
                    print(f"❌ PowerPoint extraction failed: {e}", flush=True)
                    text = ""

            # Image OCR
            elif mimetype.startswith("image/"):
                try:
                    from PIL import Image
                    import pytesseract
                    img = Image.open(tmp_path)
                    text = pytesseract.image_to_string(img)
                    print(
                        f"✅ Extracted {len(text)} chars from Image (OCR)", flush=True)
                except Exception as e:
                    print(f"❌ Image OCR failed: {e}", flush=True)
                    text = ""

            # Plain text files
            elif mimetype.startswith("text/"):
                try:
                    with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text = f.read()
                    print(
                        f"✅ Extracted {len(text)} chars from text file", flush=True)
                except Exception as e:
                    print(f"❌ Text extraction failed: {e}", flush=True)
                    text = ""

            # Store extracted content
            db2 = SessionLocal()
            try:
                # Check if content already exists
                existing = db2.query(DocumentContent).filter(
                    DocumentContent.document_id == doc_id
                ).first()

                if existing:
                    existing.content = text
                    existing.ocr_result = text
                    print(f"📝 Updated content for document {doc_id}", flush=True)
                else:
                    cc = DocumentContent(
                        document_id=doc_id,
                        content=text,
                        ocr_result=text
                    )
                    db2.add(cc)
                    print(f"📝 Created content for document {doc_id}", flush=True)

                db2.commit()
                print(
                    f"✅ Content saved to database for document {doc_id}", flush=True)

            except Exception as e:
                print(f"❌ Database error: {e}", flush=True)
                db2.rollback()
            finally:
                db2.close()

        except Exception as e:
            print(f"❌ Extraction error: {e}", flush=True)

        finally:
            # Clean up temporary file
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.unlink(tmp_path)
                except:
                    pass

    ext = file.filename.split(".")[-1]
    background_tasks.add_task(
        extract_and_store, doc.id, content_bytes, ext, file.content_type)

    print(f"✅ Document uploaded: {doc.name} (ID: {doc.id})", flush=True)
    return doc


# ========================================
# LIST DOCUMENTS (UPDATED with eager loading for new fields)
# ========================================

@router.get("/", response_model=List[DocumentOut])
def list_documents(
    directory_id: int = None,
    status: Optional[StatusEnum] = Query(StatusEnum.ACTIVE),
    file_type: Optional[str] = Query(None, description="Filter by file type"),
    category_id: Optional[int] = Query(None, description="Filter by category"),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """List documents with filtering options"""
    print(f"📄 LIST DOCUMENTS - User: {current_user.name}", flush=True)

    accessible_depts = get_accessible_departments(current_user, db)
    accessible_orgs = get_accessible_organizations(current_user, db)

    query = db.query(Document).options(
        joinedload(Document.department),
        joinedload(Document.organization),
        joinedload(Document.creator),
        joinedload(Document.document_category),
        defer(Document.data)
    ).filter(
        Document.directory_id == directory_id,
        Document.status == status
    )

    if accessible_orgs != [0]:
        query = query.filter(Document.organization_id.in_(accessible_orgs))

    if accessible_depts != [0]:
        query = query.filter(Document.department_id.in_(accessible_depts))

    if file_type:
        query = query.filter(Document.file_type == file_type)

    if category_id:
        query = query.filter(Document.document_category_id == category_id)

    documents = query.all()
    print(f"📄 Found {len(documents)} documents", flush=True)

    return documents


@router.get("/expired", response_model=List[DocumentOut])
def list_expired_documents(
    include_archived: bool = Query(
        False, description="Include already archived documents"),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """
    Get all expired documents
    Access: All users (filtered by organization)
    """
    print(f"⏰ LIST EXPIRED DOCUMENTS - User: {current_user.name}", flush=True)

    query = db.query(Document).options(
        joinedload(Document.department),
        joinedload(Document.organization),
        joinedload(Document.document_category),
        defer(Document.data)
    ).filter(
        Document.expire_date.isnot(None),
        Document.expire_date < datetime.datetime.utcnow()
    )

    if not include_archived:
        query = query.filter(Document.status == StatusEnum.ACTIVE)

    # Filter by user permissions
    if current_user.role != 'super_admin':
        query = query.filter(Document.organization_id ==
                             current_user.organization_id)

    documents = query.all()
    print(f"⏰ Found {len(documents)} expired documents", flush=True)

    return documents


@router.get("/expiring-soon", response_model=List[DocumentOut])
def list_expiring_soon_documents(
    days: int = Query(7, ge=1, le=90, description="Days threshold"),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """
    Get documents expiring within specified days
    Access: All users (filtered by organization)
    """
    print(
        f"⚠️ LIST EXPIRING SOON - User: {current_user.name}, Days: {days}", flush=True)

    threshold_date = datetime.datetime.utcnow() + timedelta(days=days)

    query = db.query(Document).options(
        joinedload(Document.department),
        joinedload(Document.organization),
        joinedload(Document.document_category),
        defer(Document.data)
    ).filter(
        Document.expire_date.isnot(None),
        Document.expire_date > datetime.datetime.utcnow(),
        Document.expire_date <= threshold_date,
        Document.status == StatusEnum.ACTIVE
    )

    # Filter by user permissions
    if current_user.role != 'super_admin':
        query = query.filter(Document.organization_id ==
                             current_user.organization_id)

    documents = query.all()
    print(f"⚠️ Found {len(documents)} documents expiring soon", flush=True)

    return documents


@router.post("/auto-archive-expired")
def auto_archive_expired(
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """
    Manually trigger auto-archive for expired documents
    Access: super_admin, org_admin only
    """
    print(f"🔄 AUTO-ARCHIVE EXPIRED - User: {current_user.name}", flush=True)

    if current_user.role not in ['super_admin', 'org_admin']:
        raise HTTPException(403, "Only admins can trigger auto-archive")

    # Build query based on role
    query = db.query(Document).filter(
        Document.expire_date < datetime.datetime.utcnow(),
        Document.status == StatusEnum.ACTIVE
    )

    # Org admin can only archive their organization's documents
    if current_user.role == 'org_admin':
        query = query.filter(Document.organization_id ==
                             current_user.organization_id)

    # Archive expired documents
    affected_count = query.update({
        Document.status: StatusEnum.ARCHIVED,
        Document.archived_at: datetime.datetime.utcnow()
    }, synchronize_session=False)

    db.commit()

    print(f"✅ Auto-archived {affected_count} expired documents", flush=True)

    return {
        "success": True,
        "message": f"Auto-archived {affected_count} expired documents",
        "affected_items": affected_count
    }


@router.get("/archived", response_model=List[DocumentOut])
def list_archived_documents(
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Get all archived documents"""
    return db.query(Document).options(
        joinedload(Document.department),
        joinedload(Document.organization),
        joinedload(Document.document_category),
        defer(Document.data)
    ).filter(Document.status == StatusEnum.ARCHIVED).all()


@router.get("/trash", response_model=List[DocumentOut])
def list_trashed_documents(
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Get all trashed documents"""
    return db.query(Document).options(
        joinedload(Document.department),
        joinedload(Document.organization),
        joinedload(Document.document_category),
        defer(Document.data)
    ).filter(Document.status == StatusEnum.TRASHED).all()


@router.get("/trash", response_model=List[DocumentOut])
def list_trashed_documents(
    db: Session = Depends(get_db)
):
    """Get all trashed documents - OPTIMIZED"""
    return db.query(Document).options(
        joinedload(Document.department),
        joinedload(Document.organization)
    ).filter(Document.status == StatusEnum.TRASHED).all()


@router.get("/search", response_model=List[DocumentOut])
def search_documents(
    q: Optional[str] = Query(None, description="Search title or content"),
    title: Optional[str] = Query(None, description="Title contains"),
    has_words: Optional[str] = Query(None, description="Content must include all these words"),
    folder_id: Optional[int] = Query(None, description="Filter by folder"),
    mimetype: Optional[str] = Query(None, description="Filter by MIME type"),
    status: Optional[StatusEnum] = Query(StatusEnum.ACTIVE, description="Filter by status"),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """
    Search documents by title and content (extracted text from files)
    """
    print(f"🔍 SEARCH - Query: {q}, Title: {title}, Has words: {has_words}", flush=True)
    
    # Base query with eager loading
    query = db.query(Document).options(
        joinedload(Document.department),
        joinedload(Document.organization),
        joinedload(Document.document_category),
        defer(Document.data)  # Don't load binary data
    ).filter(Document.status == status)
    
    # Apply permission filters
    accessible_depts = get_accessible_departments(current_user, db)
    accessible_orgs = get_accessible_organizations(current_user, db)
    
    if accessible_orgs != [0]:
        query = query.filter(Document.organization_id.in_(accessible_orgs))
    
    if accessible_depts != [0]:
        query = query.filter(Document.department_id.in_(accessible_depts))
    
    # Search in title and content
    if q:
        # Join with DocumentContent for searching extracted text
        query = query.outerjoin(DocumentContent).filter(
            or_(
                Document.title_document.ilike(f"%{q}%"),
                Document.name.ilike(f"%{q}%"),
                DocumentContent.content.ilike(f"%{q}%")
            )
        )
    
    # Search by title only
    if title:
        query = query.filter(
            or_(
                Document.title_document.ilike(f"%{title}%"),
                Document.name.ilike(f"%{title}%")
            )
        )
    
    # Search for specific words in content
    if has_words:
        words = has_words.split()
        for word in words:
            query = query.outerjoin(DocumentContent).filter(
                DocumentContent.content.ilike(f"%{word}%")
            )
    
    # Filter by folder
    if folder_id is not None:
        query = query.filter(Document.directory_id == folder_id)
    
    # Filter by mimetype
    if mimetype:
        query = query.filter(Document.mimetype == mimetype)
    
    documents = query.distinct().all()
    print(f"🔍 Found {len(documents)} documents", flush=True)
    
    return documents

@router.get("/advanced-search", response_model=List[DocumentOut])
def advanced_search_documents(
    keyword: Optional[str] = Query(None),
    file_type: Optional[str] = Query(None),
    document_category_id: Optional[int] = Query(None),
    file_category: Optional[str] = Query(None),
    file_owner: Optional[str] = Query(None),
    expire_from: Optional[str] = Query(None),
    expire_to: Optional[str] = Query(None),
    created_from: Optional[str] = Query(None),
    created_to: Optional[str] = Query(None),
    expire_status: Optional[str] = Query(None),

    # FIX: support both org_id (frontend) and organization_id (backend)
    organization_id: Optional[int] = Query(None),
    org_id: Optional[int] = Query(None),

    department_id: Optional[int] = Query(None),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """
    Corrected & professional advanced search – stable and secure.
    """

    # Normalize org_id
    if org_id and not organization_id:
        organization_id = org_id

    # Base query
    query = db.query(Document).filter(Document.status == StatusEnum.ACTIVE)

    # Authorization
    accessible_orgs = get_accessible_organizations(current_user, db)
    accessible_depts = get_accessible_departments(current_user, db)

    # 1) Limit by accessible orgs
    if accessible_orgs != [0]:
        query = query.filter(Document.organization_id.in_(accessible_orgs))

    # 2) Limit by accessible departments
    if accessible_depts != [0]:
        query = query.filter(Document.department_id.in_(accessible_depts))

    # Extra rule: non-super-admin cannot request another org
    if organization_id:
        if current_user.role != "super_admin" and organization_id not in accessible_orgs:
            raise HTTPException(status_code=403, detail="Unauthorized organization filter")

    # Apply explicit org/admin filters
    if organization_id:
        query = query.filter(Document.organization_id == organization_id)

    if department_id:
        # org_admin boleh filter dept, user hanya dept dalam accessible_depts
        if department_id not in accessible_depts and current_user.role != 'super_admin':
            raise HTTPException(403, "Unauthorized department filter")
        query = query.filter(Document.department_id == department_id)

    # Keyword
    if keyword:
        query = query.outerjoin(DocumentContent).filter(
            or_(
                Document.name.ilike(f"%{keyword}%"),
                Document.title_document.ilike(f"%{keyword}%"),
                DocumentContent.content.ilike(f"%{keyword}%")
            )
        )

    # File type
    if file_type:
        query = query.filter(Document.file_type == file_type)

    # Document category
    if document_category_id:
        query = query.filter(Document.document_category_id == document_category_id)

    # File category
    if file_category:
        query = query.filter(Document.file_category == file_category)

    # File owner
    if file_owner:
        query = query.filter(Document.file_owner.ilike(f"%{file_owner}%"))

    # Date parsing helper
    def parse_date(val: str):
        for fmt in ("%Y-%m-%d", "%Y/%m/%d", "%d-%m-%Y"):
            try:
                return datetime.strptime(val, fmt)
            except:
                pass
        raise ValueError("Invalid date format")

    # Expire from/to
    if expire_from:
        try:
            d = parse_date(expire_from)
            query = query.filter(Document.expire_date.isnot(None),
                                 Document.expire_date >= d)
        except:
            pass

    if expire_to:
        try:
            d = parse_date(expire_to) + timedelta(days=1)
            query = query.filter(Document.expire_date.isnot(None),
                                 Document.expire_date < d)
        except:
            pass

    # Expiration status
    now = datetime.datetime.utcnow()
    if expire_status == "valid":
        future7 = now + timedelta(days=7)
        query = query.filter(Document.expire_date > future7)

    elif expire_status == "expiring-soon":
        future7 = now + timedelta(days=7)
        query = query.filter(
            Document.expire_date.isnot(None),
            Document.expire_date > now,
            Document.expire_date <= future7
        )

    elif expire_status == "expired":
        query = query.filter(
            Document.expire_date.isnot(None),
            Document.expire_date < now
        )

    # Created date
    if created_from:
        try:
            d = parse_date(created_from)
            query = query.filter(Document.created_at >= d)
        except:
            pass

    if created_to:
        try:
            d = parse_date(created_to) + timedelta(days=1)
            query = query.filter(Document.created_at < d)
        except:
            pass

    # Eager loading
    query = query.options(
        joinedload(Document.organization),
        joinedload(Document.department),
        joinedload(Document.creator),
        joinedload(Document.document_category),
        defer(Document.data)
    )

    results = query.distinct().order_by(Document.created_at.desc()).all()
    return results

@router.post("/reextract-all")
async def reextract_all_documents(
    background_tasks: BackgroundTasks,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """
    Re-extract content for all documents that have empty content
    Access: super_admin, org_admin
    """
    print(f"🔄 RE-EXTRACT ALL - User: {current_user.name}", flush=True)
    
    if current_user.role not in ['super_admin', 'org_admin']:
        raise HTTPException(403, "Only admins can trigger bulk re-extraction")
    
    # Find documents with empty or missing content
    query = db.query(Document).outerjoin(DocumentContent).filter(
        or_(
            DocumentContent.content == None,
            DocumentContent.content == ""
        )
    )
    
    # Org admin only their organization
    if current_user.role == 'org_admin':
        query = query.filter(Document.organization_id == current_user.organization_id)
    
    docs = query.all()
    
    print(f"📋 Found {len(docs)} documents with empty content", flush=True)
    
    # Extract content for each document in background
    for doc in docs:
        ext = doc.name.split(".")[-1] if "." in doc.name else "txt"
        
        # Use the same extract_and_store function from upload
        def extract_and_store(doc_id: int, data: bytes, ext: str, mimetype: str):
            text = ""
            tmp_path = None
            
            try:
                with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmpf:
                    tmpf.write(data)
                    tmpf.flush()
                    tmp_path = tmpf.name
                
                # PDF extraction
                if ext.lower() == "pdf":
                    try:
                        text = extract_text(tmp_path)
                        print(f"✅ Extracted {len(text)} chars from PDF {doc_id}", flush=True)
                    except Exception as e:
                        print(f"❌ PDF extraction failed for {doc_id}: {e}", flush=True)
                
                # Word documents
                elif ext.lower() in ("docx", "doc"):
                    try:
                        d = DocxDoc(tmp_path)
                        paragraphs = [p.text for p in d.paragraphs if p.text.strip()]
                        text = "\n".join(paragraphs)
                        print(f"✅ Extracted {len(text)} chars from Word {doc_id}", flush=True)
                    except Exception as e:
                        print(f"❌ Word extraction failed for {doc_id}: {e}", flush=True)
                
                # Excel spreadsheets
                elif ext.lower() in ("xlsx", "xls"):
                    try:
                        wb = openpyxl.load_workbook(tmp_path, read_only=True, data_only=True)
                        parts = []
                        for sheet in wb.worksheets:
                            for row in sheet.iter_rows(values_only=True):
                                row_text = " ".join(str(c) for c in row if c is not None)
                                if row_text.strip():
                                    parts.append(row_text)
                        text = "\n".join(parts)
                        wb.close()
                        print(f"✅ Extracted {len(text)} chars from Excel {doc_id}", flush=True)
                    except Exception as e:
                        print(f"❌ Excel extraction failed for {doc_id}: {e}", flush=True)
                
                # PowerPoint presentations
                elif ext.lower() in ("pptx", "ppt"):
                    try:
                        prs = Presentation(tmp_path)
                        parts = []
                        for slide in prs.slides:
                            for shape in slide.shapes:
                                if hasattr(shape, "text") and shape.text.strip():
                                    parts.append(shape.text)
                        text = "\n".join(parts)
                        print(f"✅ Extracted {len(text)} chars from PowerPoint {doc_id}", flush=True)
                    except Exception as e:
                        print(f"❌ PowerPoint extraction failed for {doc_id}: {e}", flush=True)
                
                # Image OCR
                elif mimetype.startswith("image/"):
                    try:
                        img = Image.open(tmp_path)
                        text = pytesseract.image_to_string(img)
                        print(f"✅ Extracted {len(text)} chars from Image {doc_id} (OCR)", flush=True)
                    except Exception as e:
                        print(f"❌ Image OCR failed for {doc_id}: {e}", flush=True)
                
                # Store extracted content
                db2 = SessionLocal()
                try:
                    existing = db2.query(DocumentContent).filter(
                        DocumentContent.document_id == doc_id
                    ).first()
                    
                    if existing:
                        existing.content = text
                        existing.ocr_result = text
                    else:
                        cc = DocumentContent(
                            document_id=doc_id,
                            content=text,
                            ocr_result=text
                        )
                        db2.add(cc)
                    
                    db2.commit()
                    print(f"✅ Content saved for document {doc_id}", flush=True)
                    
                except Exception as e:
                    print(f"❌ Database error for {doc_id}: {e}", flush=True)
                    db2.rollback()
                finally:
                    db2.close()
            
            except Exception as e:
                print(f"❌ Extraction error for {doc_id}: {e}", flush=True)
            
            finally:
                if tmp_path and os.path.exists(tmp_path):
                    try:
                        os.unlink(tmp_path)
                    except:
                        pass
        
        background_tasks.add_task(
            extract_and_store,
            doc.id,
            doc.data,
            ext,
            doc.mimetype
        )
    
    return {
        "success": True,
        "message": f"Started re-extraction for {len(docs)} documents",
        "total_documents": len(docs)
    }

@router.post("/{document_id}/reextract")
async def reextract_document_content(
    document_id: int,
    background_tasks: BackgroundTasks,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """
    Re-extract content from document (useful if initial extraction failed)
    Access: Any user with document access
    """
    print(f"🔄 RE-EXTRACT - Document ID: {document_id}", flush=True)
    
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(404, "Document not found")
    
    if not can_access_document(current_user, document, db):
        raise HTTPException(403, "Access denied")
    
    # Extract file extension
    ext = document.name.split(".")[-1] if "." in document.name else "txt"
    
    def extract_and_store(doc_id: int, data: bytes, ext: str, mimetype: str):
        """Extract text from various file formats"""
        text = ""
        tmp_path = None
        
        try:
            # Create temporary file
            with tempfile.NamedTemporaryFile(suffix=f".{ext}", delete=False) as tmpf:
                tmpf.write(data)
                tmpf.flush()
                tmp_path = tmpf.name
            
            # PDF extraction
            if ext.lower() == "pdf":
                try:
                    from pdfminer.high_level import extract_text
                    text = extract_text(tmp_path)
                    print(f"✅ Extracted {len(text)} chars from PDF", flush=True)
                except Exception as e:
                    print(f"❌ PDF extraction failed: {e}", flush=True)
                    text = ""
            
            # Word documents
            elif ext.lower() in ("docx", "doc"):
                try:
                    from docx import Document as DocxDoc
                    d = DocxDoc(tmp_path)
                    paragraphs = [p.text for p in d.paragraphs if p.text.strip()]
                    text = "\n".join(paragraphs)
                    print(f"✅ Extracted {len(text)} chars from Word", flush=True)
                except Exception as e:
                    print(f"❌ Word extraction failed: {e}", flush=True)
                    text = ""
            
            # Excel spreadsheets
            elif ext.lower() in ("xlsx", "xls"):
                try:
                    import openpyxl
                    wb = openpyxl.load_workbook(tmp_path, read_only=True, data_only=True)
                    parts = []
                    for sheet in wb.worksheets:
                        for row in sheet.iter_rows(values_only=True):
                            row_text = " ".join(str(c) for c in row if c is not None)
                            if row_text.strip():
                                parts.append(row_text)
                    text = "\n".join(parts)
                    wb.close()
                    print(f"✅ Extracted {len(text)} chars from Excel", flush=True)
                except Exception as e:
                    print(f"❌ Excel extraction failed: {e}", flush=True)
                    text = ""
            
            # PowerPoint presentations
            elif ext.lower() in ("pptx", "ppt"):
                try:
                    from pptx import Presentation
                    prs = Presentation(tmp_path)
                    parts = []
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, "text") and shape.text.strip():
                                parts.append(shape.text)
                    text = "\n".join(parts)
                    print(f"✅ Extracted {len(text)} chars from PowerPoint", flush=True)
                except Exception as e:
                    print(f"❌ PowerPoint extraction failed: {e}", flush=True)
                    text = ""
            
            # Image OCR
            elif mimetype.startswith("image/"):
                try:
                    from PIL import Image
                    import pytesseract
                    img = Image.open(tmp_path)
                    text = pytesseract.image_to_string(img)
                    print(f"✅ Extracted {len(text)} chars from Image (OCR)", flush=True)
                except Exception as e:
                    print(f"❌ Image OCR failed: {e}", flush=True)
                    text = ""
            
            # Plain text files
            elif mimetype.startswith("text/"):
                try:
                    with open(tmp_path, 'r', encoding='utf-8', errors='ignore') as f:
                        text = f.read()
                    print(f"✅ Extracted {len(text)} chars from text file", flush=True)
                except Exception as e:
                    print(f"❌ Text extraction failed: {e}", flush=True)
                    text = ""
            
            # Store extracted content
            db2 = SessionLocal()
            try:
                # Check if content already exists
                existing = db2.query(DocumentContent).filter(
                    DocumentContent.document_id == doc_id
                ).first()
                
                if existing:
                    existing.content = text
                    existing.ocr_result = text
                    print(f"📝 Updated content for document {doc_id}", flush=True)
                else:
                    cc = DocumentContent(
                        document_id=doc_id, 
                        content=text, 
                        ocr_result=text
                    )
                    db2.add(cc)
                    print(f"📝 Created content for document {doc_id}", flush=True)
                
                db2.commit()
                print(f"✅ Content saved to database for document {doc_id}", flush=True)
                
            except Exception as e:
                print(f"❌ Database error: {e}", flush=True)
                db2.rollback()
            finally:
                db2.close()
        
        except Exception as e:
            print(f"❌ Extraction error: {e}", flush=True)
        
        finally:
            # Clean up temporary file
            if tmp_path and os.path.exists(tmp_path):
                try:
                    os.unlink(tmp_path)
                except:
                    pass
    
    background_tasks.add_task(
        extract_and_store, 
        document.id, 
        document.data, 
        ext, 
        document.mimetype
    )
    
    return {
        "success": True,
        "message": "Content extraction started in background",
        "document_id": document_id
    }

@router.put("/{document_id}/archive", response_model=StatusUpdateResponse)
def archive_document(
    document_id: int,
    db: Session = Depends(get_db)
):
    """Archive a document"""
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(404, "Document not found")

    if document.status != StatusEnum.ACTIVE:
        raise HTTPException(
            400, f"Cannot archive document with status: {document.status.value}")

    document.status = StatusEnum.ARCHIVED
    document.archived_at = datetime.datetime.utcnow()
    db.add(document)
    db.commit()

    return StatusUpdateResponse(
        success=True,
        message=f"Document '{document.name}' has been archived",
        affected_items=1
    )


@router.put("/{document_id}/restore", response_model=StatusUpdateResponse)
def restore_document(
    document_id: int,
    db: Session = Depends(get_db)
):
    """Restore a document from archive or trash"""
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(404, "Document not found")

    if document.status == StatusEnum.ACTIVE:
        raise HTTPException(400, "Document is already active")

    # Check if parent directory exists and is active (if has parent)
    if document.directory_id:
        directory = db.query(Directory).filter(
            Directory.id == document.directory_id).first()
        if not directory or directory.status != StatusEnum.ACTIVE:
            raise HTTPException(
                400, "Cannot restore: parent directory is not active")

    document.status = StatusEnum.ACTIVE
    document.archived_at = None
    document.trashed_at = None
    db.add(document)
    db.commit()

    return StatusUpdateResponse(
        success=True,
        message=f"Document '{document.name}' has been restored",
        affected_items=1
    )


@router.put("/{document_id}/trash", response_model=StatusUpdateResponse)
def move_document_to_trash(
    document_id: int,
    db: Session = Depends(get_db)
):
    """Move a document to trash (soft delete)"""
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(404, "Document not found")

    if document.status == StatusEnum.TRASHED:
        raise HTTPException(400, "Document is already in trash")

    document.status = StatusEnum.TRASHED
    document.trashed_at = datetime.datetime.utcnow()
    db.add(document)
    db.commit()

    return StatusUpdateResponse(
        success=True,
        message=f"Document '{document.name}' has been moved to trash",
        affected_items=1
    )


@router.delete("/{document_id}/permanent", response_model=StatusUpdateResponse)
def delete_document_permanent(
    document_id: int,
    db: Session = Depends(get_db)
):
    """Permanently delete a document"""
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(404, "Document not found")

    document_name = document.name
    db.delete(document)
    db.commit()

    return StatusUpdateResponse(
        success=True,
        message=f"Document '{document_name}' has been permanently deleted",
        affected_items=1
    )


@router.post("/bulk-archive", response_model=StatusUpdateResponse)
def bulk_archive_documents(
    request: BulkActionRequest,
    db: Session = Depends(get_db)
):
    """Archive multiple documents - OPTIMIZED"""
    if request.item_type != "document":
        raise HTTPException(400, "This endpoint only supports documents")

    # ✅ OPTIMIZATION: Bulk update instead of individual updates
    affected_count = db.query(Document).filter(
        Document.id.in_(request.item_ids),
        Document.status == StatusEnum.ACTIVE
    ).update({
        Document.status: StatusEnum.ARCHIVED,
        Document.archived_at: datetime.datetime.utcnow()
    }, synchronize_session=False)

    db.commit()

    return StatusUpdateResponse(
        success=True,
        message=f"Successfully archived {affected_count} documents",
        affected_items=affected_count
    )


@router.post("/bulk-trash", response_model=StatusUpdateResponse)
def bulk_trash_documents(
    request: BulkActionRequest,
    db: Session = Depends(get_db)
):
    """Move multiple documents to trash - OPTIMIZED"""
    if request.item_type != "document":
        raise HTTPException(400, "This endpoint only supports documents")

    # ✅ OPTIMIZATION: Bulk update
    affected_count = db.query(Document).filter(
        Document.id.in_(request.item_ids),
        Document.status != StatusEnum.TRASHED
    ).update({
        Document.status: StatusEnum.TRASHED,
        Document.trashed_at: datetime.datetime.utcnow()
    }, synchronize_session=False)

    db.commit()

    return StatusUpdateResponse(
        success=True,
        message=f"Successfully moved {affected_count} documents to trash",
        affected_items=affected_count
    )


@router.post("/bulk-restore", response_model=StatusUpdateResponse)
def bulk_restore_documents(
    request: BulkActionRequest,
    db: Session = Depends(get_db)
):
    """Restore multiple documents - OPTIMIZED"""
    if request.item_type != "document":
        raise HTTPException(400, "This endpoint only supports documents")

    # Get documents with their directories
    documents = db.query(Document).filter(
        Document.id.in_(request.item_ids),
        Document.status != StatusEnum.ACTIVE
    ).all()

    affected_count = 0
    valid_doc_ids = []

    for doc in documents:
        # Check parent directory status
        if doc.directory_id:
            directory = db.query(Directory).filter(
                Directory.id == doc.directory_id).first()
            if not directory or directory.status != StatusEnum.ACTIVE:
                continue
        valid_doc_ids.append(doc.id)
        affected_count += 1

    # Bulk update valid documents
    if valid_doc_ids:
        db.query(Document).filter(Document.id.in_(valid_doc_ids)).update({
            Document.status: StatusEnum.ACTIVE,
            Document.archived_at: None,
            Document.trashed_at: None
        }, synchronize_session=False)

    db.commit()

    return StatusUpdateResponse(
        success=True,
        message=f"Successfully restored {affected_count} documents",
        affected_items=affected_count
    )


@router.delete("/bulk-permanent", response_model=StatusUpdateResponse)
def bulk_delete_documents_permanent(
    request: BulkActionRequest,
    db: Session = Depends(get_db)
):
    """Permanently delete multiple documents - OPTIMIZED"""
    if request.item_type != "document":
        raise HTTPException(400, "This endpoint only supports documents")

    # ✅ OPTIMIZATION: Bulk delete
    affected_count = db.query(Document).filter(
        Document.id.in_(request.item_ids)
    ).delete(synchronize_session=False)

    db.commit()

    return StatusUpdateResponse(
        success=True,
        message=f"Successfully deleted {affected_count} documents permanently",
        affected_items=affected_count
    )


# ===== PREVIEW & DOWNLOAD ENDPOINTS (unchanged) =====

@router.get("/{document_id}/preview")
async def preview_document(
    document_id: int,
    page: int = Query(1, ge=1, description="Page number for PDF"),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(404, "Document not found")

    if not can_access_document(current_user, document, db):
        raise HTTPException(403, "Access denied")

    mimetype = document.mimetype.lower()

    if mimetype == "application/pdf":
        try:
            import fitz
            pdf_stream = BytesIO(document.data)
            pdf_document = fitz.open(stream=pdf_stream, filetype="pdf")

            total_pages = len(pdf_document)
            if page > total_pages:
                pdf_document.close()
                raise HTTPException(
                    400, f"Page {page} exceeds document length {total_pages}")

            pdf_page = pdf_document[page - 1]
            mat = fitz.Matrix(2, 2)
            pix = pdf_page.get_pixmap(matrix=mat)
            img_data = pix.tobytes("png")

            pdf_document.close()
            pdf_stream.close()

            return Response(
                content=img_data,
                media_type="image/png",
                headers={
                    "X-Total-Pages": str(total_pages),
                    "X-Current-Page": str(page),
                    "Cache-Control": "max-age=3600"
                }
            )
        except ImportError:
            raise HTTPException(500, "PyMuPDF not installed")
        except Exception as e:
            raise HTTPException(500, f"Error processing PDF: {str(e)}")

    elif mimetype.startswith("image/"):
        return Response(
            content=document.data,
            media_type=document.mimetype,
            headers={"Cache-Control": "max-age=3600"}
        )

    elif mimetype in [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword"
    ]:
        try:
            from docx import Document as DocxDoc
            import tempfile

            with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp_docx:
                tmp_docx.write(document.data)
                tmp_docx_path = tmp_docx.name

            try:
                docx_doc = DocxDoc(tmp_docx_path)
                html_parts = [
                    '<!DOCTYPE html><html><head><meta charset="UTF-8">']
                html_parts.append(
                    '<style>body{font-family:Arial,sans-serif;padding:20px;max-width:800px;margin:0 auto;}</style>')
                html_parts.append('</head><body>')

                for i, para in enumerate(docx_doc.paragraphs[:100]):
                    if para.text.strip():
                        html_parts.append(f'<p>{para.text}</p>')
                    if i >= 99:
                        html_parts.append(
                            '<p><em>... (content truncated for preview)</em></p>')
                        break

                html_parts.append('</body></html>')
                html_content = ''.join(html_parts)

                return Response(
                    content=html_content,
                    media_type="text/html",
                    headers={"Cache-Control": "max-age=3600"}
                )
            finally:
                if os.path.exists(tmp_docx_path):
                    os.remove(tmp_docx_path)
        except Exception as e:
            raise HTTPException(
                500, f"Error processing Word document: {str(e)}")
    else:
        raise HTTPException(400, f"Preview not supported for {mimetype}")


@router.get("/{document_id}/download")
async def download_document(
    document_id: int,
    db: Session = Depends(get_db)
):
    """Download the original document"""
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(404, "Document not found")

    return Response(
        content=document.data,
        media_type=document.mimetype,
        headers={"Content-Disposition": f'attachment; filename="{document.name}"'}
    )


@router.get("/{document_id}/info")
async def get_document_info(
    document_id: int,
    db: Session = Depends(get_db)
):
    """Get document information with all details"""
    document = db.query(Document).options(
        joinedload(Document.department),
        joinedload(Document.organization),
        joinedload(Document.document_category),
        joinedload(Document.creator)
    ).filter(Document.id == document_id).first()
    
    if not document:
        raise HTTPException(404, "Document not found")

    info = {
        "id": document.id,
        "name": document.name,
        "title": document.title_document,
        "mimetype": document.mimetype,
        "size": document.size,
        "total_pages": document.total_pages or 1,
        "created_at": document.created_at,
        
        # NEW FIELDS
        "file_type": document.file_type,
        "file_category": document.file_category,
        "file_owner": document.file_owner,
        "expire_date": document.expire_date,
        
        # Relations
        "document_category": {
            "id": document.document_category.id,
            "name": document.document_category.name,
            "code": document.document_category.code
        } if document.document_category else None,
        
        "organization": {
            "id": document.organization.id,
            "name": document.organization.name,
            "code": document.organization.code
        } if document.organization else None,
        
        "department": {
            "id": document.department.id,
            "name": document.department.name,
            "code": document.department.code
        } if document.department else None
    }

    # Get PDF page count if applicable
    if document.mimetype == "application/pdf":
        try:
            import fitz
            pdf_document = fitz.open(stream=BytesIO(document.data), filetype="pdf")
            info["total_pages"] = len(pdf_document)
            pdf_document.close()
        except:
            pass

    return info

# ===== SHARE ENDPOINTS (unchanged but with fix) =====

SHARE_CACHE_DIR = "temp/share_tokens"
os.makedirs(SHARE_CACHE_DIR, exist_ok=True)


@router.post("/{document_id}/share")
def create_share_link(
    document_id: int,
    expires_in: int = Query(7),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Create a simple share link"""
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(404, "Document not found")

    # SPECIAL RULE: User can share documents they own regardless of org/dept
    if document.created_by == current_user.id:
        # Owner can always share their own documents
        pass  # Allow access
    else:
        # Check organization and department access for non-owners
        accessible_depts = get_accessible_departments(current_user, db)
        accessible_orgs = get_accessible_organizations(current_user, db)

        has_access = (
            document.organization_id in accessible_orgs and
            document.department_id in accessible_depts
        )

        if not has_access:
            raise HTTPException(403, "You don't have access to this document")

    share_token = secrets.token_urlsafe(32)
    expires_at = datetime.datetime.now(timezone.utc) + timedelta(days=expires_in)
    created_at = datetime.datetime.now(timezone.utc)

    share_data = {
        "document_id": document_id,
        "share_token": share_token,
        "expires_at": expires_at.isoformat(),
        "created_by": current_user.id,
        "created_at": created_at.isoformat(),
        "document_name": document.name,
        "document_title": document.title_document
    }

    os.makedirs(SHARE_CACHE_DIR, exist_ok=True)
    token_file = os.path.join(SHARE_CACHE_DIR, f"{share_token}.json")
    with open(token_file, "w") as f:
        json.dump(share_data, f)

    share_url = f"/shared/{share_token}"

    return {
        "share_link": share_url,
        "share_token": share_token,
        "expires_at": expires_at.isoformat(),
        "expires_in_days": expires_in
    }


# Continue with remaining share endpoints...
# (Include all remaining endpoints from original file)

@router.get("/shared/{share_token}/preview")
def preview_shared_document(
    share_token: str,
    page: int = Query(1, ge=1, description="Page number for PDF"),
):
    """Preview shared document - public endpoint"""
    token_file = os.path.join(SHARE_CACHE_DIR, f"{share_token}.json")

    if not os.path.exists(token_file):
        raise HTTPException(404, "Share link not found or has been deleted")

    try:
        with open(token_file, "r") as f:
            share_data = json.load(f)
    except Exception as e:
        raise HTTPException(500, f"Error reading share data: {str(e)}")

    # PERBAIKAN: Gunakan datetime.now(timezone.utc)
    expires_at = datetime.datetime.fromisoformat(share_data["expires_at"])
    if datetime.datetime.now(timezone.utc) > expires_at:
        try:
            os.remove(token_file)
        except:
            pass
        raise HTTPException(403, "Share link has expired")

    # Get document
    db = SessionLocal()
    try:
        document_id = share_data["document_id"]
        document = db.query(Document).filter(
            Document.id == document_id).first()

        if not document:
            raise HTTPException(404, "Document not found")

        # Handle PDF preview
        mimetype = document.mimetype.lower()

        if mimetype == "application/pdf":
            try:
                import fitz

                pdf_stream = BytesIO(document.data)
                pdf_document = fitz.open(stream=pdf_stream, filetype="pdf")

                total_pages = len(pdf_document)
                if page > total_pages:
                    pdf_document.close()
                    raise HTTPException(
                        400, f"Page {page} exceeds document length {total_pages}")

                pdf_page = pdf_document[page - 1]
                mat = fitz.Matrix(2, 2)
                pix = pdf_page.get_pixmap(matrix=mat)

                img_data = pix.tobytes("png")

                pdf_document.close()
                pdf_stream.close()

                return Response(
                    content=img_data,
                    media_type="image/png",
                    headers={
                        "X-Total-Pages": str(total_pages),
                        "X-Current-Page": str(page),
                        "Cache-Control": "max-age=3600"
                    }
                )

            except ImportError:
                raise HTTPException(500, "PyMuPDF not installed")
            except Exception as e:
                raise HTTPException(500, f"Error processing PDF: {str(e)}")

        # Handle Images
        elif mimetype.startswith("image/"):
            return Response(
                content=document.data,
                media_type=document.mimetype,
                headers={"Cache-Control": "max-age=3600"}
            )

        # Handle Word documents
        elif mimetype in [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword"
        ]:
            try:
                from docx import Document as DocxDoc
                import tempfile

                with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp_docx:
                    tmp_docx.write(document.data)
                    tmp_docx_path = tmp_docx.name

                try:
                    docx_doc = DocxDoc(tmp_docx_path)

                    html_parts = [
                        '<!DOCTYPE html><html><head><meta charset="UTF-8">']
                    html_parts.append(
                        '<style>body{font-family:Arial,sans-serif;padding:20px;max-width:800px;margin:0 auto;}</style>')
                    html_parts.append('</head><body>')

                    for i, para in enumerate(docx_doc.paragraphs[:100]):
                        if para.text.strip():
                            html_parts.append(f'<p>{para.text}</p>')
                        if i >= 99:
                            html_parts.append(
                                '<p><em>... (content truncated)</em></p>')
                            break

                    html_parts.append('</body></html>')
                    html_content = ''.join(html_parts)

                    return Response(
                        content=html_content,
                        media_type="text/html",
                        headers={"Cache-Control": "max-age=3600"}
                    )
                finally:
                    if os.path.exists(tmp_docx_path):
                        os.remove(tmp_docx_path)

            except Exception as e:
                raise HTTPException(
                    500, f"Error processing Word document: {str(e)}")

        else:
            raise HTTPException(400, f"Preview not supported for {mimetype}")

    finally:
        db.close()


@router.get("/shared/{share_token}/info")
def get_shared_document_info(share_token: str):
    """Get shared document info"""
    token_file = os.path.join(SHARE_CACHE_DIR, f"{share_token}.json")

    if not os.path.exists(token_file):
        raise HTTPException(404, "Share link not found")

    try:
        with open(token_file, "r") as f:
            share_data = json.load(f)
    except:
        raise HTTPException(500, "Error reading share data")

    # Check expiration
    expires_at = datetime.datetime.fromisoformat(share_data["expires_at"])
    # if datetime.datetime.datetime.datetime.utcnow() > expires_at:
    #     try:
    #         os.remove(token_file)
    #     except:
    #         pass
    #     raise HTTPException(403, "Share link has expired")

    # Get document
    db = SessionLocal()
    try:
        document_id = share_data["document_id"]
        document = db.query(Document).filter(
            Document.id == document_id).first()

        if not document:
            raise HTTPException(404, "Document not found")

        info = {
            "id": document.id,
            "name": document.name,
            "title": document.title_document,
            "mimetype": document.mimetype,
            "size": document.size,
            "total_pages": 1,
            "created_at": document.created_at,
            "shared_by": share_data.get("created_by", "Unknown"),
            "shared_at": share_data.get("created_at"),
            "expires_at": share_data.get("expires_at")
        }

        # Get page count for PDFs
        if document.mimetype == "application/pdf":
            try:
                import fitz
                pdf_document = fitz.open(
                    stream=BytesIO(document.data), filetype="pdf")
                info["total_pages"] = len(pdf_document)
                pdf_document.close()
            except:
                pass

        return info
    finally:
        db.close()


@router.get("/shared/{share_token}/download")
def download_shared_document(share_token: str):
    """Download shared document"""
    token_file = os.path.join(SHARE_CACHE_DIR, f"{share_token}.json")

    if not os.path.exists(token_file):
        raise HTTPException(404, "Share link not found")

    try:
        with open(token_file, "r") as f:
            share_data = json.load(f)
    except:
        raise HTTPException(500, "Error reading share data")

    # Check expiration
    expires_at = datetime.fromisoformat(share_data["expires_at"])
    if datetime.datetime.datetime.utcnow() > expires_at:
        try:
            os.remove(token_file)
        except:
            pass
        raise HTTPException(403, "Share link has expired")

    # Get document
    db = SessionLocal()
    try:
        document_id = share_data["document_id"]
        document = db.query(Document).filter(
            Document.id == document_id).first()

        if not document:
            raise HTTPException(404, "Document not found")

        return Response(
            content=document.data,
            media_type=document.mimetype,
            headers={
                "Content-Disposition": f'attachment; filename="{document.name}"'
            }
        )
    finally:
        db.close()


@router.delete("/{share_token}/revoke")
def revoke_share_link(share_token: str):
    """Revoke a share link"""
    token_file = os.path.join(SHARE_CACHE_DIR, f"{share_token}.json")

    if not os.path.exists(token_file):
        raise HTTPException(404, "Share link not found")

    try:
        os.remove(token_file)
        return {
            "success": True,
            "message": "Share link revoked successfully"
        }
    except Exception as e:
        raise HTTPException(500, f"Failed to revoke share link: {str(e)}")


@router.get("/shared/with-me", response_model=List[DocumentOut])
def get_shared_with_me(
    skip: int = 0,
    limit: int = 50,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    shared_docs = db.query(Document).join(DocumentShare).filter(
        DocumentShare.expires_at > datetime.datetime.datetime.utcnow(),
        (
            (DocumentShare.target_user_id == current_user.id) |
            (DocumentShare.target_department_id == current_user.department_id) |
            (DocumentShare.target_organization_id == current_user.organization_id)
        )
    ).offset(skip).limit(limit).all()

    return shared_docs


@router.get("/{document_id}/edit")
async def get_editable_document(
    document_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Get document data for editing"""
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(404, "Document not found")

    if not can_access_document(current_user, document, db):
        raise HTTPException(403, "Access denied")

    # Return document data for editing
    return {
        "id": document.id,
        "name": document.name,
        "title": document.title_document,
        "mimetype": document.mimetype,
        "data": document.data.decode('latin-1') if document.data else "",
        "size": document.size
    }


@router.put("/{document_id}/edit")
async def update_document_content(
    document_id: int,
    content: str = Form(...),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Update document content"""
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(404, "Document not found")

    if not can_access_document(current_user, document, db):
        raise HTTPException(403, "Access denied")

    # Update document data
    document.data = content.encode('latin-1')
    document.size = len(content)
    document.updated_at = datetime.datetime.datetime.utcnow()

    db.add(document)
    db.commit()
    db.refresh(document)

    return {"message": "Document updated successfully"}


@router.get("/{document_id}/preview/excel")
async def preview_excel(
    document_id: int,
    sheet: int = Query(0, description="Sheet index"),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Preview Excel file as HTML table"""
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(404, "Document not found")

    if not document.mimetype in [
        "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        "application/vnd.ms-excel"
    ]:
        raise HTTPException(400, "Not an Excel file")

    try:
        # Create temporary file
        with tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False) as tmp:
            tmp.write(document.data)
            tmp_path = tmp.name

        # Read Excel file
        df = pd.read_excel(tmp_path, sheet_name=sheet)

        # Convert to HTML
        html_content = df.to_html(classes="excel-preview-table", index=False)

        # Clean up
        os.unlink(tmp_path)

        # Return HTML response
        full_html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <style>
                .excel-preview-table {{
                    border-collapse: collapse;
                    width: 100%;
                    font-family: Arial, sans-serif;
                }}
                .excel-preview-table th, .excel-preview-table td {{
                    border: 1px solid #ddd;
                    padding: 8px;
                    text-align: left;
                }}
                .excel-preview-table th {{
                    background-color: #f2f2f2;
                    font-weight: bold;
                }}
                .excel-preview-table tr:nth-child(even) {{
                    background-color: #f9f9f9;
                }}
            </style>
        </head>
        <body>
            <div class="excel-preview">
                <h3>Excel Preview - Sheet {sheet + 1}</h3>
                {html_content}
            </div>
        </body>
        </html>
        """

        return Response(content=full_html, media_type="text/html")

    except Exception as e:
        raise HTTPException(500, f"Error processing Excel file: {str(e)}")


@router.get("/{document_id}/preview/word")
async def preview_word(
    document_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Preview Word document as HTML"""
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(404, "Document not found")

    if not document.mimetype in [
        "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        "application/msword"
    ]:
        raise HTTPException(400, "Not a Word document")

    try:
        # Create temporary file
        with tempfile.NamedTemporaryFile(suffix=".docx", delete=False) as tmp:
            tmp.write(document.data)
            tmp_path = tmp.name

        # Read Word document
        doc = DocxDoc(tmp_path)

        # Extract text content
        full_text = []
        for para in doc.paragraphs:
            if para.text.strip():
                full_text.append(f"<p>{para.text}</p>")

        # Extract tables
        for table in doc.tables:
            table_html = "<table class='word-table'>"
            for row in table.rows:
                table_html += "<tr>"
                for cell in row.cells:
                    table_html += f"<td>{cell.text}</td>"
                table_html += "</tr>"
            table_html += "</table>"
            full_text.append(table_html)

        # Clean up
        os.unlink(tmp_path)

        # Return HTML response
        full_html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <style>
                body {{
                    font-family: Arial, sans-serif;
                    line-height: 1.6;
                    margin: 20px;
                    background: #f5f5f5;
                }}
                .word-preview {{
                    background: white;
                    padding: 30px;
                    border-radius: 8px;
                    box-shadow: 0 2px 10px rgba(0,0,0,0.1);
                    max-width: 800px;
                    margin: 0 auto;
                }}
                .word-table {{
                    border-collapse: collapse;
                    width: 100%;
                    margin: 15px 0;
                }}
                .word-table td {{
                    border: 1px solid #ddd;
                    padding: 8px;
                    text-align: left;
                }}
                .word-table tr:nth-child(even) {{
                    background-color: #f9f9f9;
                }}
            </style>
        </head>
        <body>
            <div class="word-preview">
                <h3>Word Document Preview</h3>
                {''.join(full_text)}
            </div>
        </body>
        </html>
        """

        return Response(content=full_html, media_type="text/html")

    except Exception as e:
        raise HTTPException(500, f"Error processing Word document: {str(e)}")


@router.get("/{document_id}/preview/ppt")
async def preview_ppt(
    document_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Preview PowerPoint file with better formatting"""
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(404, "Document not found")

    if not document.mimetype in [
        "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        "application/vnd.ms-powerpoint"
    ]:
        raise HTTPException(400, "Not a PowerPoint file")

    try:
        # Create temporary file
        with tempfile.NamedTemporaryFile(suffix=".pptx", delete=False) as tmp:
            tmp.write(document.data)
            tmp_path = tmp.name

        # Read PowerPoint file
        prs = Presentation(tmp_path)

        # Generate HTML with better slide formatting
        slides_html = []
        for i, slide in enumerate(prs.slides):
            slide_content = []

            # Add slide title if exists
            title = ""
            if slide.shapes.title:
                title = slide.shapes.title.text

            # Process all shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    # Check if it's a title
                    if shape == slide.shapes.title:
                        slide_content.append(f"<h4>{shape.text}</h4>")
                    else:
                        # Check formatting for bullet points
                        text = shape.text.strip()
                        if any(char in text for char in ['•', '-', '*']):
                            lines = text.split('\n')
                            bullet_list = "<ul>"
                            for line in lines:
                                if line.strip():
                                    bullet_list += f"<li>{line.strip().lstrip('•-* ')}</li>"
                            bullet_list += "</ul>"
                            slide_content.append(bullet_list)
                        else:
                            slide_content.append(f"<p>{text}</p>")

            slides_html.append(f"""
            <div class="slide">
                <div class="slide-header">
                    <span class="slide-number">Slide {i + 1}</span>
                    {f'<h3 class="slide-title">{title}</h3>' if title else ''}
                </div>
                <div class="slide-content">
                    {''.join(slide_content)}
                </div>
            </div>
            """)

        # Clean up
        os.unlink(tmp_path)

        # Return HTML response
        full_html = f"""
        <!DOCTYPE html>
        <html>
        <head>
            <style>
                body {{
                    font-family: 'Segoe UI', Tahoma, Geneva, Verdana, sans-serif;
                    background: linear-gradient(135deg, #667eea 0%, #764ba2 100%);
                    margin: 0;
                    padding: 20px;
                    min-height: 100vh;
                }}
                .ppt-preview {{
                    max-width: 1000px;
                    margin: 0 auto;
                }}
                .ppt-preview h2 {{
                    color: white;
                    text-align: center;
                    margin-bottom: 30px;
                    text-shadow: 0 2px 4px rgba(0,0,0,0.3);
                }}
                .slide {{
                    background: white;
                    border-radius: 15px;
                    margin: 25px 0;
                    padding: 30px;
                    box-shadow: 0 10px 30px rgba(0,0,0,0.2);
                    border: 1px solid rgba(255,255,255,0.3);
                }}
                .slide-header {{
                    border-bottom: 3px solid #667eea;
                    padding-bottom: 15px;
                    margin-bottom: 20px;
                    display: flex;
                    justify-content: space-between;
                    align-items: center;
                }}
                .slide-number {{
                    background: #667eea;
                    color: white;
                    padding: 5px 15px;
                    border-radius: 20px;
                    font-weight: bold;
                    font-size: 14px;
                }}
                .slide-title {{
                    color: #333;
                    margin: 0;
                    flex-grow: 1;
                    text-align: center;
                }}
                .slide-content h4 {{
                    color: #667eea;
                    margin: 15px 0 10px 0;
                    border-left: 4px solid #667eea;
                    padding-left: 10px;
                }}
                .slide-content p {{
                    margin: 10px 0;
                    line-height: 1.6;
                    color: #555;
                }}
                .slide-content ul {{
                    margin: 10px 0;
                    padding-left: 20px;
                }}
                .slide-content li {{
                    margin: 5px 0;
                    line-height: 1.5;
                }}
            </style>
        </head>
        <body>
            <div class="ppt-preview">
                <h2>PowerPoint Presentation - {len(slides_html)} Slides</h2>
                {''.join(slides_html)}
            </div>
        </body>
        </html>
        """

        return Response(content=full_html, media_type="text/html")

    except Exception as e:
        raise HTTPException(500, f"Error processing PowerPoint file: {str(e)}")


@router.get("/{document_id}/content/text")
async def get_document_text_content(
    document_id: int,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Get readable text content from document for editing"""
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(404, "Document not found")

    try:
        text_content = ""

        # Create temporary file
        with tempfile.NamedTemporaryFile(suffix=f".{document.name.split('.')[-1]}", delete=False) as tmp:
            tmp.write(document.data)
            tmp_path = tmp.name

        # Extract text based on file type
        if document.mimetype in [
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            "application/msword"
        ]:
            # Word document
            doc = DocxDoc(tmp_path)
            text_content = "\n".join(
                [para.text for para in doc.paragraphs if para.text.strip()])

        elif document.mimetype in [
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            "application/vnd.ms-excel"
        ]:
            # Excel document - extract from first sheet
            df = pd.read_excel(tmp_path, sheet_name=0)
            text_content = df.to_string(index=False)

        elif document.mimetype in [
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            "application/vnd.ms-powerpoint"
        ]:
            # PowerPoint document
            prs = Presentation(tmp_path)
            for i, slide in enumerate(prs.slides):
                text_content += f"=== Slide {i+1} ===\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += f"{shape.text}\n"
                text_content += "\n"

        else:
            # Try to extract text using existing method
            doc_content = db.query(DocumentContent).filter(
                DocumentContent.document_id == document_id).first()
            if doc_content and doc_content.content:
                text_content = doc_content.content
            else:
                text_content = "No readable content found"

        # Clean up
        os.unlink(tmp_path)

        return {
            "id": document.id,
            "name": document.name,
            "title": document.title_document,
            "mimetype": document.mimetype,
            "content": text_content,
            "size": document.size
        }

    except Exception as e:
        raise HTTPException(
            500, f"Error extracting document content: {str(e)}")


@router.put("/{document_id}/content/text")
async def update_document_text_content(
    document_id: int,
    content_update: dict,
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """Update document text content (for simple text-based editing)"""
    document = db.query(Document).filter(Document.id == document_id).first()
    if not document:
        raise HTTPException(404, "Document not found")

    # For now, we'll update the DocumentContent table
    # In a real implementation, you'd want to update the actual Office file
    doc_content = db.query(DocumentContent).filter(
        DocumentContent.document_id == document_id).first()

    if not doc_content:
        doc_content = DocumentContent(document_id=document_id)

    doc_content.content = content_update.get("content", "")
    doc_content.ocr_result = content_update.get("content", "")

    db.add(doc_content)
    db.commit()

    return {"message": "Document content updated successfully"}

@router.get("/", response_model=List[DocumentCategoryOut])
def list_document_categories(
    org_id: Optional[int] = Query(None, description="Filter by organization id (optional)"),
    db: Session = Depends(get_db),
    current_user: User = Depends(get_current_user)
):
    """
    Return document categories. If caller is not super_admin and org_id not provided,
    only return categories for current_user.organization_id.
    """
    q = db.query(DocumentCategory)

    # if caller supplies org_id AND is super_admin -> allow arbitrary org_id
    if org_id:
        if current_user.role != "super_admin":
            # Non-super-admin may not query other orgs
            if org_id != current_user.organization_id:
                raise HTTPException(403, "Cannot query categories for other organizations")
        q = q.filter(DocumentCategory.organization_id == org_id)
    else:
        # if org_id not passed -> for non-super admins limit to user's org
        if current_user.role != "super_admin":
            q = q.filter(DocumentCategory.organization_id == current_user.organization_id)

    return q.order_by(DocumentCategory.name).all()
